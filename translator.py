#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Unified Office Document Translator
---------------------------------
Translates text content in Microsoft Office documents (Excel, Word, PowerPoint)
between Japanese, English, and Vietnamese using the Gemini API.

This script combines the functionality of trans-excel.py and trans-office.py
into a single, unified interface with improved error handling, progress tracking,
and a modern command-line interface.
"""

import os
import sys
import time
import argparse
import json
import re
import glob
import subprocess
import mimetypes
from enum import Enum
from typing import List, Dict, Tuple, Optional, Union, Callable, Any
from pathlib import Path
import datetime

# Add rich for improved terminal UI
try:
    from rich.console import Console
    from rich.progress import Progress, SpinnerColumn, TextColumn, BarColumn, TaskID
    from rich.panel import Panel
    from rich.text import Text
    HAS_RICH = True
except ImportError:
    HAS_RICH = False

# Enhanced PowerPoint processing imports
try:
    import zipfile
    import xml.etree.ElementTree as ET
    from lxml import etree
    HAS_LXML = True
except ImportError:
    HAS_LXML = False

# Image processing for enhanced PowerPoint support
try:
    import base64
    import io
    HAS_IMAGE_LIBS = True
except ImportError:
    HAS_IMAGE_LIBS = False

try:
    import win32com.client
    import comtypes
    HAS_COM = True
except ImportError:
    HAS_COM = False

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

# Define document type enum
class DocType(Enum):
    """Enum for supported document types"""
    EXCEL = "excel"
    WORD = "word"
    POWERPOINT = "powerpoint"
    UNKNOWN = "unknown"

# Initialize rich console if available
console = Console() if HAS_RICH else None

def print_info(message: str) -> None:
    """Print an informational message with rich formatting if available"""
    if HAS_RICH:
        console.print(f"[bold blue]ℹ️[/bold blue] {message}")
    else:
        print(f"ℹ️ {message}")

def print_success(message: str) -> None:
    """Print a success message with rich formatting if available"""
    if HAS_RICH:
        console.print(f"[bold green]✅[/bold green] {message}")
    else:
        print(f"✅ {message}")

def print_warning(message: str) -> None:
    """Print a warning message with rich formatting if available"""
    if HAS_RICH:
        console.print(f"[bold yellow]⚠️[/bold yellow] {message}")
    else:
        print(f"⚠️ {message}")

def print_error(message: str) -> None:
    """Print an error message with rich formatting if available"""
    if HAS_RICH:
        console.print(f"[bold red]❌[/bold red] {message}")
    else:
        print(f"❌ {message}")

def print_header(title: str) -> None:
    """Print a header with rich formatting if available"""
    if HAS_RICH:
        console.print(Panel(Text(title, justify="center"), style="bold blue"))
    else:
        border = "=" * (len(title) + 10)
        print(f"\n{border}\n    {title}\n{border}")

def check_and_install_dependencies() -> bool:
    """
    Check if required dependencies are installed and install if needed.
    
    Returns:
        bool: True if all dependencies are available, False otherwise
    """
    try:
        # Create requirements file
        script_dir = os.path.dirname(os.path.abspath(__file__))
        req_file = os.path.join(script_dir, "translator-requirements.txt")
        
        if not os.path.exists(req_file):
            print_info("Creating requirements file...")
            with open(req_file, 'w', encoding='utf-8') as f:
                f.write(
                    "openai>=1.0.0\n"
                    "xlwings>=0.30.0\n"
                    "python-pptx>=0.6.21\n"
                    "python-docx>=0.8.11\n"
                    "python-dotenv>=1.0.0\n"
                    "pathlib>=1.0.1\n"
                    "rich>=13.0.0\n"
                    "tqdm>=4.66.0\n"
                )
            print_success(f"Requirements file created at: {req_file}")
        
        print_info("Checking for required libraries...")
        
        # Check if packages are installed
        missing_packages = []
        required_packages = {
            "openai": "OpenAI API client",
            "xlwings": "Excel file processing",
            "pptx": "PowerPoint file processing (python-pptx)",
            "docx": "Word file processing (python-docx)",
            "dotenv": "Environment variable loading (python-dotenv)",
            "rich": "Enhanced terminal output"
        }
        
        for package, description in required_packages.items():
            try:
                __import__(package)
            except ImportError:
                missing_packages.append(package)
                print_warning(f"Missing {description} ({package})")
        
        # If packages are missing, attempt to install them
        if missing_packages:
            print_info(f"Missing packages: {', '.join(missing_packages)}")
            print_info("Attempting to install missing packages automatically...")
            
            try:
                # Install required packages
                subprocess.check_call([sys.executable, "-m", "pip", "install", "-r", req_file])
                print_success("Successfully installed required packages.")
            except subprocess.CalledProcessError as e:
                print_error(f"Failed to install packages automatically: {str(e)}")
                print_info("Please manually install the required packages with this command:")
                print_info(f"pip install -r {req_file}")
                return False
        else:
            print_success("All required libraries are already installed.")
            
        # Verify imports after installation attempt
        try:
            import openai
            import xlwings as xw
            import pptx
            import docx
            from dotenv import load_dotenv
            
            # Try to import rich again if it was missing
            if "rich" in missing_packages:
                import rich
                from rich.console import Console
                from rich.progress import Progress, SpinnerColumn, TextColumn, BarColumn
                from rich.panel import Panel
                from rich.text import Text
                global console, HAS_RICH
                console = Console()
                HAS_RICH = True
                print_success("Rich library loaded successfully - enhanced UI enabled.")
                
            print_success("All required libraries loaded successfully.")
            return True
        except ImportError as e:
            print_error(f"Error importing library after installation attempt: {str(e)}")
            print_info("Please manually install the required libraries and try again:")
            print_info(f"pip install -r {req_file}")
            return False
    except Exception as e:
        print_error(f"Error checking libraries: {str(e)}")
        return False

def get_file_type(file_path: str) -> DocType:
    """
    Determine the document type based on file extension.
    
    Args:
        file_path: Path to the file
        
    Returns:
        DocType: Type of document (EXCEL, WORD, POWERPOINT, or UNKNOWN)
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext in ['.xlsx', '.xls', '.xlsm']:
        return DocType.EXCEL
    elif ext in ['.docx', '.doc']:
        return DocType.WORD
    elif ext in ['.pptx', '.ppt']:
        return DocType.POWERPOINT
    else:
        return DocType.UNKNOWN

def clean_text(text: Optional[str]) -> str:
    """
    Clean and normalize text before translation.
    
    Args:
        text: The text to clean
        
    Returns:
        str: Cleaned text
    """
    if not text or not isinstance(text, str):
        return ""
    text = ' '.join(text.split())  # Normalize whitespace
    return text.strip()

def should_translate(text: Optional[str]) -> bool:
    """
    Check if text needs translation.
    
    Args:
        text: The text to check
        
    Returns:
        bool: True if text should be translated, False otherwise
    """
    text = clean_text(text)
    if not text or len(text) < 2:
        return False
    if re.match(r'^[\d\s,.-]+$', text):  # Contains only numbers and number formatting characters
        return False
    if isinstance(text, str) and text.startswith('='):  # Excel formula
        return False
    return True

def translate_batch(
    texts: List[str], 
    target_lang: str = "ja", 
    api_client: Any = None,
    progress_callback: Optional[Callable[[int, int], None]] = None
) -> List[str]:
    """
    Translate a batch of texts to the target language.
    
    Args:
        texts: List of texts to translate
        target_lang: Target language code ("ja" for Japanese, "vi" for Vietnamese, "en" for English)
        api_client: OpenAI client instance
        progress_callback: Optional callback to report progress
        
    Returns:
        List[str]: List of translated texts
    """
    if not texts:
        return []

    # Read system prompt from file
    script_dir = os.path.dirname(os.path.abspath(__file__))
    prompt_file = os.path.join(script_dir, "translator-system-prompt.txt")
    
    # Check if the prompt file exists
    if os.path.exists(prompt_file):
        with open(prompt_file, 'r', encoding='utf-8') as f:
            system_prompt = f.read()
    else:
        # Use default prompt if file doesn't exist
        system_prompt = """You are a professional translator. Follow these rules strictly:
1. Output ONLY the translation, nothing else
2. DO NOT include the original text in your response
3. DO NOT add any explanations or notes
4. Keep IDs, model numbers, and special characters unchanged
5. Use standard terminology for technical terms
6. Preserve the original formatting (spaces, line breaks)
7. Use proper grammar and punctuation
8. Only keep unchanged: proper names, IDs, and technical codes
9. Translate all segments separated by "|||" and keep them separated with the same delimiter"""
        # Create default prompt file
        with open(prompt_file, 'w', encoding='utf-8') as f:
            f.write(system_prompt)
        print_info(f"Default prompt file created at: {prompt_file}")

    # Combine texts with separator
    separator = "|||"
    combined_text = separator.join(texts)

    # Determine translation direction based on parameter
    language_map = {
        "ja": "to Japanese",
        "vi": "to Vietnamese", 
        "en": "to English",
        "th": "to Thai",
        "zh": "to Chinese (Simplified)",
        "ko": "to Korean"
    }
    direction = language_map.get(target_lang, "to the target language")
    
    user_prompt = f"Translate the following text {direction}, keeping segments separated by '{separator}':\n\n{combined_text}"

    # Debug: Show first part of what we're translating
    if len(texts) > 0:
        print_info(f"Translating {len(texts)} texts {direction}")
        first_text = texts[0][:50] + "..." if len(texts[0]) > 50 else texts[0]
        print_info(f"First text sample: {first_text}")

    # Add retry mechanism for API calls
    retries = 0
    max_retries = 3  # Maximum number of retries for API calls
    delay_seconds = 2  # Delay between API calls
    
    while retries <= max_retries:
        try:
            # Update progress if callback provided
            if progress_callback:
                progress_callback(retries, max_retries)
                
            # Call translation API
            print_info(f"Calling translation API (attempt {retries+1}/{max_retries+1})...")
            response = api_client.chat.completions.create(
                model="gemini-2.0-flash",  # Adjust model as needed
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ]
            )
            
            # Debug response
            print_info("API response received, analyzing structure...")
            
            # Extract content from response (handle different response structures)
            translated_text = None
            
            # Try accessing as a dictionary first (most reliable method for unknown structures)
            try:
                response_dict = {}
                if hasattr(response, '__dict__'):
                    response_dict = response.__dict__
                elif hasattr(response, 'model_dump'):
                    response_dict = response.model_dump()
                else:
                    import json
                    try:
                        response_dict = json.loads(json.dumps(response, default=lambda o: o.__dict__))
                    except:
                        # Last resort, try string conversion
                        response_str = str(response)
                        if HAS_RICH:
                            print_info("Raw response preview:")
                            console.print(response_str[:100] + "...", style="dim")
                        else:
                            print_info(f"Raw response: {response_str[:100]}...")
                
                # Try to debug response structure on first attempt
                if retries == 0 and HAS_RICH:
                    keys = list(response_dict.keys()) if isinstance(response_dict, dict) else 'Not a dictionary'
                    print_info(f"Response keys: {keys}")
            except Exception as dict_err:
                print_warning(f"Could not convert response to dictionary: {str(dict_err)}")
            
            # First try standard OpenAI format
            try:
                if hasattr(response, 'choices') and response.choices and len(response.choices) > 0:
                    if hasattr(response.choices[0], 'message') and hasattr(response.choices[0].message, 'content'):
                        translated_text = response.choices[0].message.content
                        print_success("Found content in standard OpenAI format")
            except Exception as std_err:
                print_warning(f"Could not extract via standard path: {str(std_err)}")
            
            # If still no content, try alternative paths or direct access
            if not translated_text:
                try:
                    # Try direct string conversion (some APIs return the text directly)
                    if hasattr(response, '__str__'):
                        response_str = str(response).strip()
                        if response_str and len(response_str) > 0 and not response_str.startswith('{'):
                            translated_text = response_str
                            print_success("Using direct string conversion")
                except Exception as str_err:
                    print_warning(f"String conversion failed: {str(str_err)}")
            
            # Last resort: look for any attribute that might contain our text
            if not translated_text:
                print_info("Searching for content in response attributes...")
                content_attributes = ['content', 'text', 'answer', 'result', 'output', 'response', 'data']
                
                # Check direct attributes first
                for attr in content_attributes:
                    if hasattr(response, attr):
                        potential_content = getattr(response, attr)
                        if isinstance(potential_content, str) and potential_content.strip():
                            translated_text = potential_content
                            print_success(f"Found content in '{attr}' attribute")
                            break
                
                # If still not found, check nested attributes
                if not translated_text and isinstance(response_dict, dict):
                    # Check all dictionary paths that might contain the content
                    for attr in content_attributes:
                        if attr in response_dict and isinstance(response_dict[attr], str) and response_dict[attr].strip():
                            translated_text = response_dict[attr]
                            print_success(f"Found content in response_dict['{attr}']")
                            break
            
            # If we still don't have content, raise error with details
            if not translated_text:
                response_repr = str(response)[:500]  # Get truncated string representation
                raise ValueError(f"Could not extract translation content from API response of type {type(response).__name__}. Response preview: {response_repr}")

            # Split translation result into separate parts
            translated_parts = translated_text.split(separator)
            
            # Debug: Show first translated result
            if len(translated_parts) > 0:
                first_translated = translated_parts[0][:50] + "..." if len(translated_parts[0]) > 50 else translated_parts[0]
                print_info(f"First translation result: {first_translated}")

            # Handle case when number of translated parts doesn't match
            if len(translated_parts) != len(texts):
                print_warning(f"Number of translated parts ({len(translated_parts)}) doesn't match number of original texts ({len(texts)})")
                # Ensure number of translated parts equals number of original texts
                if len(translated_parts) < len(texts):
                    translated_parts.extend(texts[len(translated_parts):])
                else:
                    translated_parts = translated_parts[:len(texts)]

            # Delay to avoid exceeding API limits
            time.sleep(delay_seconds)
            return translated_parts

        except Exception as e:
            retries += 1
            if retries > max_retries:
                print_error(f"Error translating batch after {max_retries} retries: {str(e)}")
                # Return original texts if translation fails
                return texts
            
            print_warning(f"API call failed (attempt {retries}/{max_retries}): {str(e)}")
            print_info(f"Retrying in {retries * 2} seconds...")
            time.sleep(retries * 2)  # Exponential backoff
            
    # This point should not be reached due to the return in the exception handler above
    return texts 

def process_excel_file(
    input_path: str, 
    target_lang: str = "ja", 
    api_client: Any = None,
    progress: Optional[Any] = None
) -> Optional[str]:
    """
    Process Excel file: read, translate and save with original format.
    
    Args:
        input_path: Path to the Excel file
        target_lang: Target language code
        api_client: OpenAI client instance
        progress: Optional rich progress instance
        
    Returns:
        Optional[str]: Path to the translated file, or None if processing failed
    """
    try:
        # Import xlwings here to ensure it's loaded
        import xlwings as xw
        
        # Create output file path
        filename = os.path.basename(input_path)
        base_name, ext = os.path.splitext(filename)

        # Create output directory at the same level as the script
        project_dir = os.path.dirname(os.path.abspath(__file__))
        output_dir = os.path.join(project_dir, "output")
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f"{base_name}-translated{ext}")

        print_info(f"Processing Excel file: {filename}")

        # Task IDs for progress tracking
        scan_task = None
        file_task = None
        
        if progress:
            file_task = progress.add_task(f"[cyan]Processing {filename}", total=1.0)
            scan_task = progress.add_task("[green]Scanning content", total=1.0)

        # Open workbook with xlwings to preserve formatting
        app = xw.App(visible=False)
        wb = None  # Initialize wb
        try:
            wb = app.books.open(input_path)
            
            if progress:
                # Update progress: file opened
                progress.update(file_task, completed=0.2, description=f"[cyan]Processing {filename} - File opened")
                # Update scanning task: starting scan
                progress.update(scan_task, completed=0.1, description="[green]Scanning sheets")

            # Loop through each sheet
            sheet_count = len(wb.sheets)
            total_cells_to_translate = 0
            total_cells_translated = 0
            
            # First pass: count all elements to translate for better progress tracking
            sheet_translation_data = []
            
            for sheet_idx, sheet in enumerate(wb.sheets):
                if progress:
                    progress.update(scan_task, completed=0.1 + (0.4 * (sheet_idx / sheet_count)), 
                                    description=f"[green]Scanning sheet {sheet_idx+1}/{sheet_count}")
                
                # Collect data from cells that need translation
                texts_to_translate = []
                cell_references = []

                # Scan through used data range
                used_rng = sheet.used_range
                if used_rng.count > 1 or used_rng.value is not None:  # Only scan if sheet has data
                    for cell in used_rng:
                        # Check if cell.value is None before calling str()
                        cell_value_str = str(cell.value) if cell.value is not None else ""
                        if cell_value_str and should_translate(cell_value_str):
                            texts_to_translate.append(clean_text(cell_value_str))
                            cell_references.append(cell)
                else:
                    print_info(f"Sheet '{sheet.name}' is empty or has no data.")

                # Process shapes with text
                try:
                    shapes_collection = sheet.api.Shapes
                    shapes_count = shapes_collection.Count

                    if shapes_count > 0:
                        print_info(f"Sheet '{sheet.name}' has {shapes_count} shapes to check")

                        # Process each shape by index (Excel COM API indexes from 1)
                        for i in range(1, shapes_count + 1):
                            shape = None  # Initialize to avoid errors if .Item(i) fails
                            try:
                                shape = shapes_collection.Item(i)
                                shape_text = None

                                # --- Try multiple methods to get text from shape ---
                                
                                # Method 1: TextFrame
                                try:
                                    if hasattr(shape, 'TextFrame'):
                                        if shape.TextFrame.HasText:
                                            shape_text = shape.TextFrame.Characters().Text
                                except:
                                    pass
                                
                                # Method 2: TextFrame2
                                if not shape_text:
                                    try:
                                        if hasattr(shape, 'TextFrame2'):
                                            shape_text = shape.TextFrame2.TextRange.Text
                                    except:
                                        pass
                                
                                # Method 3: AlternativeText
                                if not shape_text:
                                    try:
                                        if hasattr(shape, 'AlternativeText') and shape.AlternativeText:
                                            shape_text = shape.AlternativeText
                                    except:
                                        pass
                                
                                # Method 4: OLEFormat (for OLE objects)
                                if not shape_text:
                                    try:
                                        if hasattr(shape, 'OLEFormat') and hasattr(shape.OLEFormat, 'Object'):
                                            if hasattr(shape.OLEFormat.Object, 'Text'):
                                                shape_text = shape.OLEFormat.Object.Text
                                    except:
                                        pass
                                
                                # Method 5: TextEffect (for WordArt)
                                if not shape_text:
                                    try:
                                        if hasattr(shape, 'TextEffect') and hasattr(shape.TextEffect, 'Text'):
                                            shape_text = shape.TextEffect.Text
                                    except:
                                        pass
                                
                                # If text is found, add to translation list
                                if shape_text and should_translate(shape_text):
                                    clean_shape_text = clean_text(shape_text)
                                    preview = clean_shape_text[:30] + "..." if len(clean_shape_text) > 30 else clean_shape_text
                                    print_info(f"Shape {i}: Found text: {preview}")
                                    texts_to_translate.append(clean_shape_text)
                                    
                                    # Save tuple with information for later updates
                                    cell_references.append(('shape', sheet, i))

                            except Exception as outer_e:
                                # General error when processing shape
                                print_warning(f"Error processing shape {i}: {str(outer_e)}")
                                continue

                except Exception as e:
                    print_warning(f"Error processing shapes on sheet '{sheet.name}': {str(e)}")

                # Save the sheet data for translation phase
                sheet_translation_data.append((sheet, texts_to_translate, cell_references))
                total_cells_to_translate += len(texts_to_translate)
            
            # Finish scanning phase
            if progress:
                progress.update(scan_task, completed=1.0, description="[green]Scanning complete")
                
                # Create a new task for translation
                trans_task = progress.add_task("[yellow]Translating content", total=total_cells_to_translate)
                # Update file task
                progress.update(file_task, completed=0.4, description=f"[cyan]Processing {filename} - Translation in progress")
            
            # Skip if no text to translate
            if total_cells_to_translate == 0:
                print_success("No text to translate in the Excel file.")
                
                if progress:
                    progress.update(file_task, completed=1.0, description=f"[cyan]Processing {filename} - No translation needed")
                
                # Save file with original format
                wb.save(output_path)
                print_success(f"File saved at: {output_path}")
                return output_path

            # Second pass: translate and update cell content
            batch_size = 50  # Smaller batch size for better progress tracking
            
            for sheet_idx, (sheet, texts_to_translate, cell_references) in enumerate(sheet_translation_data):
                if not texts_to_translate:
                    continue  # Skip sheets with no text to translate
                    
                print_info(f"Processing sheet: {sheet.name}")
                
                # Split into batches for processing
                total_batches = (len(texts_to_translate) - 1) // batch_size + 1
                
                for i in range(0, len(texts_to_translate), batch_size):
                    batch_texts = texts_to_translate[i:i+batch_size]
                    batch_refs = cell_references[i:i+batch_size]
                    current_batch_num = i // batch_size + 1

                    print_info(f"Translating batch {current_batch_num}/{total_batches} ({len(batch_texts)} texts)")
                    
                    # Define a progress callback for translation
                    def update_translation_progress(current: int, total: int) -> None:
                        if progress:
                            progress.update(trans_task, 
                                           description=f"[yellow]Translating batch {current_batch_num}/{total_batches} (Attempt {current+1}/{total+1})")
                    
                    # Translate batch
                    translated_batch = translate_batch(
                        batch_texts, 
                        target_lang,
                        api_client, 
                        progress_callback=update_translation_progress
                    )
                    
                    # Update translated content
                    print_info(f"Updating content for batch {current_batch_num}...")
                    
                    for j, ref in enumerate(batch_refs):
                        # Update progress
                        if progress:
                            total_cells_translated += 1
                            progress.update(trans_task, completed=total_cells_translated)
                            
                        # Check if index j is within translated_batch
                        if j < len(translated_batch) and translated_batch[j] is not None:
                            try:
                                # Update content for shape and cell
                                if isinstance(ref, tuple) and ref[0] == 'shape':
                                    # Process shape: ref is ('shape', sheet_obj, shape_index)
                                    _, sheet_obj, shape_index = ref  # Unpack tuple
                                    try:
                                        # Get shape object again
                                        shape_to_update = sheet_obj.api.Shapes.Item(shape_index)
                                        updated = False
                                        
                                        # Try multiple methods to update text for shape
                                        
                                        # Method 1: TextFrame
                                        try:
                                            if hasattr(shape_to_update, 'TextFrame') and shape_to_update.TextFrame.HasText:
                                                shape_to_update.TextFrame.Characters().Text = translated_batch[j]
                                                updated = True
                                        except:
                                            pass
                                            
                                        # Method 2: TextFrame2
                                        if not updated:
                                            try:
                                                if hasattr(shape_to_update, 'TextFrame2'):
                                                    shape_to_update.TextFrame2.TextRange.Text = translated_batch[j]
                                                    updated = True
                                            except:
                                                pass
                                                
                                        # Method 3: AlternativeText
                                        if not updated:
                                            try:
                                                if hasattr(shape_to_update, 'AlternativeText'):
                                                    shape_to_update.AlternativeText = translated_batch[j]
                                                    updated = True
                                            except:
                                                pass
                                                
                                        # Method 4: TextEffect (for WordArt)
                                        if not updated:
                                            try:
                                                if hasattr(shape_to_update, 'TextEffect') and hasattr(shape_to_update.TextEffect, 'Text'):
                                                    shape_to_update.TextEffect.Text = translated_batch[j]
                                                    updated = True
                                            except:
                                                pass
                                                
                                        # Method 5: OLEFormat
                                        if not updated:
                                            try:
                                                if hasattr(shape_to_update, 'OLEFormat') and hasattr(shape_to_update.OLEFormat, 'Object'):
                                                    if hasattr(shape_to_update.OLEFormat.Object, 'Text'):
                                                        shape_to_update.OLEFormat.Object.Text = translated_batch[j]
                                                        updated = True
                                            except:
                                                pass
                                                
                                        if updated:
                                            print_success(f"Updated text for shape {shape_index} on sheet '{sheet_obj.name}'")
                                        else:
                                            print_warning(f"Could not update text for shape {shape_index} on sheet '{sheet_obj.name}'")
                                        
                                    except Exception as update_err:
                                        print_warning(f"Error updating shape {shape_index} on sheet '{sheet_obj.name}': {str(update_err)}")
                                elif hasattr(ref, 'value'):  # Is a cell
                                    ref.value = translated_batch[j]
                                else:
                                    print_warning(f"Unknown reference type: {type(ref)}")

                            except Exception as update_single_err:
                                # Catch general errors when updating a specific cell/shape
                                ref_info = f"Shape index {ref[2]} on sheet {ref[1].name}" if isinstance(ref, tuple) else f"Cell {ref.address}"
                                print_warning(f"Could not update content for {ref_info}: {str(update_single_err)}")

            # Update file task
            if progress:
                progress.update(file_task, completed=0.8, description=f"[cyan]Processing {filename} - Saving file")
                progress.update(trans_task, completed=total_cells_to_translate, description="[yellow]Translation complete")

            # Save file with original format
            print_info(f"Saving translated file to: {output_path}")
            wb.save(output_path)
            print_success(f"File saved successfully: {output_path}")
            
            if progress:
                progress.update(file_task, completed=1.0, description=f"[cyan]Processing {filename} - Complete")

            return output_path

        except Exception as wb_process_err:
            print_error(f"Error processing workbook '{filename}': {str(wb_process_err)}")
            # Ensure workbook is closed if error occurs before saving
            if wb is not None:
                try:
                    wb.close()
                except Exception as close_err:
                    print_warning(f"Error trying to close workbook after processing error: {close_err}")
            if progress and file_task:
                progress.update(file_task, completed=1.0, description=f"[red]Processing {filename} - Failed")
            return None
        finally:
            # Close workbook (if not already closed) and Excel app
            # wb.close() has been called in the except block if needed
            # Just need to ensure app is closed
            if 'app' in locals() and app.pid:  # Check if app exists and is still running
                app.quit()
                print_info("Excel application closed.")

    except Exception as e:
        print_error(f"Critical error when processing Excel file '{input_path}': {str(e)}")
        if progress and file_task:
            progress.update(file_task, completed=1.0, description=f"[red]Processing {filename} - Failed")
        return None 

class AdvancedPowerPointProcessor:
    """
    Advanced PowerPoint processor with multi-engine approach to handle complex elements
    like SmartArt, WordArt, OLE objects, and complex shapes.
    """
    
    def __init__(self):
        self.engines = {
            'hybrid': self._process_with_hybrid_approach,
            'python_pptx': self._process_with_python_pptx,
            'com_automation': self._process_with_com_automation,
            'xml_direct': self._process_with_xml_direct
        }
        self.extracted_elements = []
        
    def process_presentation(self, file_path: str, target_lang: str = "ja", api_client: Any = None, progress: Optional[Any] = None) -> Optional[str]:
        """Main processing method with multiple fallback engines"""
        
        filename = os.path.basename(file_path)
        print_info(f"Starting advanced PowerPoint processing for: {filename}")
        
        # Try engines in order of capability
        for engine_name, engine_func in self.engines.items():
            try:
                print_info(f"Attempting processing with {engine_name} engine...")
                result = engine_func(file_path, target_lang, api_client, progress)
                if result:
                    print_success(f"Successfully processed with {engine_name} engine")
                    return result
            except Exception as e:
                print_warning(f"{engine_name} engine failed: {str(e)}")
                continue
                
        raise Exception("All processing engines failed")

    def _process_with_hybrid_approach(self, file_path: str, target_lang: str, api_client: Any, progress: Optional[Any]) -> Optional[str]:
        """Most comprehensive approach combining all methods"""
        
        filename = os.path.basename(file_path)
        print_info(f"Using hybrid approach for advanced PowerPoint processing")
        
        # Step 1: Extract with enhanced python-pptx (baseline)
        basic_elements = self._extract_with_enhanced_pptx(file_path, progress)
        
        # Step 2: Extract complex elements with COM automation (Windows only)
        if HAS_COM and os.name == 'nt':
            try:
                print_info("Attempting COM automation for advanced elements...")
                complex_elements = self._extract_with_com_automation_only(file_path)
                basic_elements.extend(complex_elements)
                print_info(f"COM automation found {len(complex_elements)} additional elements")
            except Exception as e:
                print_warning(f"COM automation not available: {e}")
                
        # Step 3: Extract missed elements with direct XML parsing
        if HAS_LXML:
            try:
                print_info("Attempting XML parsing for additional elements...")
                xml_elements = self._extract_with_xml_parsing(file_path)
                basic_elements.extend(xml_elements)
                print_info(f"XML parsing found {len(xml_elements)} additional elements")
            except Exception as e:
                print_warning(f"XML parsing failed: {e}")
                
        # Step 4: Process translations
        if basic_elements:
            print_info(f"Processing {len(basic_elements)} total extracted elements")
            translated_elements = self._translate_elements(basic_elements, target_lang, api_client, progress)
            return self._apply_translations(file_path, translated_elements, progress)
            
        print_warning("No elements extracted by any engine")
        return None

    def _extract_with_enhanced_pptx(self, file_path: str, progress: Optional[Any]) -> List[Dict]:
        """Enhanced python-pptx extraction with additional element types"""
        elements = []
        
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            slide_count = len(prs.slides)
            print_info(f"Extracting from {slide_count} slides using enhanced python-pptx")
            
            for slide_idx, slide in enumerate(prs.slides):
                if progress:
                    progress.update(progress.task_ids[0] if progress.task_ids else None, 
                                   description=f"[green]Enhanced scanning slide {slide_idx+1}/{slide_count}")
                
                # Standard shapes with enhanced detection
                for shape_idx, shape in enumerate(slide.shapes):
                    elements.extend(self._extract_from_shape_enhanced(shape, slide_idx, shape_idx))
                    
                # Layout and master slide elements
                try:
                    layout_elements = self._extract_from_layout(slide.slide_layout, slide_idx)
                    elements.extend(layout_elements)
                except Exception as e:
                    print_warning(f"Layout extraction failed: {e}")
                    
                # Notes with enhanced extraction
                try:
                    if hasattr(slide, "notes_slide") and slide.notes_slide:
                        notes_elements = self._extract_notes_enhanced(slide.notes_slide, slide_idx)
                        elements.extend(notes_elements)
                except Exception as e:
                    print_warning(f"Notes extraction failed: {e}")
                    
        except Exception as e:
            print_error(f"Enhanced python-pptx extraction failed: {e}")
            
        print_info(f"Enhanced python-pptx extracted {len(elements)} elements")
        return elements

    def _extract_from_shape_enhanced(self, shape, slide_idx: int, shape_idx: int) -> List[Dict]:
        """Enhanced shape text extraction with better element detection"""
        elements = []
        
        try:
            # Regular text frames with multiple paragraph support
            if hasattr(shape, "text_frame") and shape.text_frame:
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    if paragraph.text and paragraph.text.strip():
                        elements.append({
                            'type': 'text_frame_paragraph',
                            'location': (slide_idx, shape_idx, para_idx),
                            'text': paragraph.text.strip(),
                            'shape': shape,
                            'paragraph': paragraph
                        })
            
            # Enhanced table processing
            if hasattr(shape, "table"):
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if cell.text and cell.text.strip():
                            elements.append({
                                'type': 'table_cell',
                                'location': (slide_idx, shape_idx, row_idx, col_idx),
                                'text': cell.text.strip(),
                                'cell': cell
                            })
            
            # Chart elements (enhanced)
            if hasattr(shape, "chart"):
                chart_elements = self._extract_from_chart_enhanced(shape.chart, slide_idx, shape_idx)
                elements.extend(chart_elements)
                
            # Group items with recursive processing
            if hasattr(shape, "shapes"):  # Group shape
                for group_idx, group_shape in enumerate(shape.shapes):
                    group_elements = self._extract_from_shape_enhanced(group_shape, slide_idx, f"{shape_idx}-{group_idx}")
                    elements.extend(group_elements)
                    
            # Placeholder detection (improved error handling)
            if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                try:
                    # Safely check if it's a real placeholder with valid type
                    placeholder_type = getattr(shape.placeholder_format, 'type', None)
                    if placeholder_type is not None and hasattr(shape, 'text') and shape.text and shape.text.strip():
                        elements.append({
                            'type': 'placeholder',
                            'location': (slide_idx, shape_idx),
                            'text': shape.text.strip(),
                            'shape': shape,
                            'placeholder_type': placeholder_type
                        })
                except AttributeError as e:
                    # Shape has placeholder_format but accessing type failed
                    if "placeholder" in str(e).lower():
                        # Still try to extract text if available
                        if hasattr(shape, 'text') and shape.text and shape.text.strip():
                            elements.append({
                                'type': 'text_shape',
                                'location': (slide_idx, shape_idx),
                                'text': shape.text.strip(),
                                'shape': shape
                            })
                except Exception:
                    # Any other error, just skip this placeholder processing
                    pass
                    
        except Exception as e:
            print_warning(f"Error extracting from shape {shape_idx}: {e}")
            
        return elements

    def _extract_from_chart_enhanced(self, chart, slide_idx: int, shape_idx: int) -> List[Dict]:
        """Enhanced chart text extraction"""
        elements = []
        
        try:
            # Chart title
            if hasattr(chart, 'chart_title') and chart.chart_title and hasattr(chart.chart_title, 'text_frame'):
                title_text = chart.chart_title.text_frame.text
                if title_text and title_text.strip():
                    elements.append({
                        'type': 'chart_title',
                        'location': (slide_idx, shape_idx, 'title'),
                        'text': title_text.strip(),
                        'chart': chart
                    })
            
            # Axis titles
            try:
                if hasattr(chart, 'value_axis') and chart.value_axis.axis_title:
                    axis_text = chart.value_axis.axis_title.text_frame.text
                    if axis_text and axis_text.strip():
                        elements.append({
                            'type': 'chart_axis_title',
                            'location': (slide_idx, shape_idx, 'value_axis'),
                            'text': axis_text.strip(),
                            'chart': chart
                        })
            except:
                pass
                
            try:
                if hasattr(chart, 'category_axis') and chart.category_axis.axis_title:
                    axis_text = chart.category_axis.axis_title.text_frame.text
                    if axis_text and axis_text.strip():
                        elements.append({
                            'type': 'chart_axis_title',
                            'location': (slide_idx, shape_idx, 'category_axis'),
                            'text': axis_text.strip(),
                            'chart': chart
                        })
            except:
                pass
                
        except Exception as e:
            print_warning(f"Chart extraction failed: {e}")
            
        return elements

    def _extract_from_layout(self, layout, slide_idx: int) -> List[Dict]:
        """Extract text from slide layout and master slide"""
        elements = []
        
        try:
            # Layout placeholders
            for placeholder in layout.placeholders:
                if hasattr(placeholder, 'text_frame') and placeholder.text_frame.text:
                    text = placeholder.text_frame.text.strip()
                    if text and should_translate(text):
                        elements.append({
                            'type': 'layout_placeholder',
                            'location': (slide_idx, 'layout', placeholder.placeholder_format.idx),
                            'text': text,
                            'placeholder': placeholder
                        })
                        
        except Exception as e:
            print_warning(f"Layout extraction failed: {e}")
            
        return elements

    def _extract_notes_enhanced(self, notes_slide, slide_idx: int) -> List[Dict]:
        """Enhanced notes extraction"""
        elements = []
        
        try:
            if hasattr(notes_slide, "notes_text_frame") and notes_slide.notes_text_frame:
                # Extract paragraph by paragraph
                for para_idx, paragraph in enumerate(notes_slide.notes_text_frame.paragraphs):
                    if paragraph.text and paragraph.text.strip() and should_translate(paragraph.text):
                        elements.append({
                            'type': 'notes_paragraph',
                            'location': (slide_idx, 'notes', para_idx),
                            'text': paragraph.text.strip(),
                            'paragraph': paragraph
                        })
                        
        except Exception as e:
            print_warning(f"Notes extraction failed: {e}")
            
        return elements

    def _extract_with_com_automation_only(self, file_path: str) -> List[Dict]:
        """Windows COM automation for advanced elements only"""
        elements = []
        
        if not HAS_COM or os.name != 'nt':
            return elements
            
        try:
            print_info("Starting COM automation for advanced elements...")
            
            # Start PowerPoint application
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            ppt_app.Visible = False
            
            # Open presentation
            presentation = ppt_app.Presentations.Open(os.path.abspath(file_path))
            
            for slide_idx in range(presentation.Slides.Count):
                slide = presentation.Slides[slide_idx + 1]  # COM is 1-indexed
                
                for shape_idx in range(slide.Shapes.Count):
                    shape = slide.Shapes[shape_idx + 1]  # COM is 1-indexed
                    
                    # SmartArt extraction
                    if shape.Type == 15:  # SmartArt
                        smartart_elements = self._extract_smartart_com(shape, slide_idx, shape_idx)
                        elements.extend(smartart_elements)
                    
                    # WordArt extraction
                    try:
                        if hasattr(shape, 'TextEffect') and shape.TextEffect.Text:
                            wordart_text = shape.TextEffect.Text.strip()
                            if wordart_text and should_translate(wordart_text):
                                elements.append({
                                    'type': 'wordart',
                                    'location': (slide_idx, shape_idx),
                                    'text': wordart_text,
                                    'com_shape': shape,
                                    'requires_com': True
                                })
                    except:
                        pass
                    
                    # OLE Objects
                    if shape.Type == 7:  # OLE Object
                        try:
                            ole_text = self._extract_ole_object_text(shape)
                            if ole_text and should_translate(ole_text):
                                elements.append({
                                    'type': 'ole_object',
                                    'location': (slide_idx, shape_idx),
                                    'text': ole_text,
                                    'com_shape': shape,
                                    'requires_com': True
                                })
                        except:
                            pass
            
            # Close presentation and application
            presentation.Close()
            ppt_app.Quit()
            
            print_info(f"COM automation extracted {len(elements)} advanced elements")
            
        except Exception as e:
            print_warning(f"COM automation failed: {e}")
            
        return elements

    def _extract_smartart_com(self, shape, slide_idx: int, shape_idx: int) -> List[Dict]:
        """Extract text from SmartArt using COM"""
        elements = []
        
        try:
            smartart = shape.SmartArt
            
            # Extract from all nodes
            for node_idx in range(smartart.AllNodes.Count):
                node = smartart.AllNodes[node_idx + 1]  # COM is 1-indexed
                
                if hasattr(node, 'TextFrame2') and node.TextFrame2.HasText:
                    node_text = node.TextFrame2.TextRange.Text.strip()
                    if node_text and should_translate(node_text):
                        elements.append({
                            'type': 'smartart_node',
                            'location': (slide_idx, shape_idx, node_idx),
                            'text': node_text,
                            'com_shape': shape,
                            'node_index': node_idx,
                            'requires_com': True
                        })
                        
        except Exception as e:
            print_warning(f"SmartArt extraction failed: {e}")
            
        return elements

    def _extract_ole_object_text(self, shape) -> Optional[str]:
        """Extract text from OLE objects"""
        try:
            # Try different methods to get OLE object text
            if hasattr(shape, 'OLEFormat') and shape.OLEFormat:
                ole_object = shape.OLEFormat.Object
                
                # For Excel objects
                if hasattr(ole_object, 'Cells'):
                    # This would require more complex Excel COM automation
                    return "Excel Object (requires manual translation)"
                    
                # For Word objects
                if hasattr(ole_object, 'Content'):
                    return ole_object.Content.Text
                    
                # Generic text extraction
                if hasattr(ole_object, 'Text'):
                    return ole_object.Text
                    
        except:
            pass
            
        return None

    def _extract_with_xml_parsing(self, file_path: str) -> List[Dict]:
        """Direct XML parsing for elements missed by other methods"""
        elements = []
        
        if not HAS_LXML:
            return elements
            
        try:
            print_info("Starting XML parsing for additional elements...")
            
            # PowerPoint files are ZIP archives
            with zipfile.ZipFile(file_path, 'r') as pptx_zip:
                # Parse slide XML files
                slide_files = [f for f in pptx_zip.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
                
                for slide_file in slide_files:
                    slide_idx = int(slide_file.split('slide')[1].split('.')[0]) - 1
                    
                    with pptx_zip.open(slide_file) as xml_file:
                        xml_content = xml_file.read()
                        xml_elements = self._parse_slide_xml(xml_content, slide_idx)
                        elements.extend(xml_elements)
                        
            print_info(f"XML parsing extracted {len(elements)} additional elements")
            
        except Exception as e:
            print_warning(f"XML parsing failed: {e}")
            
        return elements

    def _parse_slide_xml(self, xml_content: bytes, slide_idx: int) -> List[Dict]:
        """Parse slide XML for advanced elements"""
        elements = []
        
        try:
            root = etree.fromstring(xml_content)
            
            # Define namespaces
            namespaces = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram'
            }
            
            # Find all text elements not caught by python-pptx
            text_elements = root.xpath('.//a:t[not(ancestor::p:txBody)]', namespaces=namespaces)
            
            for idx, text_elem in enumerate(text_elements):
                if text_elem.text and text_elem.text.strip() and should_translate(text_elem.text):
                    elements.append({
                        'type': 'xml_text',
                        'location': (slide_idx, f'xml_{idx}'),
                        'text': text_elem.text.strip(),
                        'xml_element': text_elem,
                        'requires_xml': True
                    })
                    
        except Exception as e:
            print_warning(f"XML slide parsing failed: {e}")
            
        return elements

    def _process_with_python_pptx(self, file_path: str, target_lang: str, api_client: Any, progress: Optional[Any]) -> Optional[str]:
        """Fallback to basic python-pptx processing"""
        print_info("Falling back to basic python-pptx processing")
        return process_powerpoint_file_basic(file_path, target_lang, api_client, progress)

    def _process_with_com_automation(self, file_path: str, target_lang: str, api_client: Any, progress: Optional[Any]) -> Optional[str]:
        """Pure COM automation approach (Windows only)"""
        if not HAS_COM or os.name != 'nt':
            raise Exception("COM automation not available on this platform")
            
        print_info("Using pure COM automation approach")
        # This would be a complete COM implementation - for now, fall back
        return self._process_with_hybrid_approach(file_path, target_lang, api_client, progress)

    def _process_with_xml_direct(self, file_path: str, target_lang: str, api_client: Any, progress: Optional[Any]) -> Optional[str]:
        """Direct XML manipulation approach"""
        if not HAS_LXML:
            raise Exception("LXML not available for XML processing")
            
        print_info("Using direct XML manipulation approach")
        # This would be a complete XML implementation - for now, fall back
        return self._process_with_hybrid_approach(file_path, target_lang, api_client, progress)

    def _translate_elements(self, elements: List[Dict], target_lang: str, api_client: Any, progress: Optional[Any]) -> List[Dict]:
        """Translate extracted elements using existing translation function"""
        if not elements:
            return []
            
        print_info(f"Translating {len(elements)} extracted elements")
        
        # Group texts for batch translation
        texts_to_translate = []
        element_map = {}
        
        for idx, element in enumerate(elements):
            if element['text'] and should_translate(element['text']):
                clean_text_value = clean_text(element['text'])
                texts_to_translate.append(clean_text_value)
                element_map[len(texts_to_translate) - 1] = idx
        
        if not texts_to_translate:
            return elements
            
        print_info(f"Processing {len(texts_to_translate)} texts for translation")
        
        # Use existing translation function
        translated_texts = translate_batch(texts_to_translate, target_lang, api_client)
        
        # Map translations back to elements
        for trans_idx, translated_text in enumerate(translated_texts):
            if trans_idx in element_map:
                original_idx = element_map[trans_idx]
                elements[original_idx]['translated_text'] = translated_text
                
        return elements

    def _apply_translations(self, file_path: str, translated_elements: List[Dict], progress: Optional[Any]) -> Optional[str]:
        """Apply translations back to the presentation"""
        
        # Create output path
        filename = os.path.basename(file_path)
        base_name, ext = os.path.splitext(filename)
        # Create completely safe filename with timestamp
        import re
        import datetime
        # Remove all special characters and replace with underscores
        safe_base_name = re.sub(r'[^\w\-]', '_', base_name).strip('_')
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        project_dir = os.path.dirname(os.path.abspath(__file__))
        output_dir = os.path.join(project_dir, "output")
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f"TRANSLATED_{safe_base_name}_{timestamp}{ext}")
        
        try:
            print_info(f"Applying translations to {len(translated_elements)} elements")
            
            # Load presentation
            from pptx import Presentation
            prs = Presentation(file_path)
            
            # Track updates
            updated_count = 0
            com_updates_needed = []
            
            # Apply translations
            for element in translated_elements:
                if 'translated_text' not in element:
                    continue
                    
                try:
                    element_type = element['type']
                    translated_text = element['translated_text']
                    original_text = element.get('text', '')[:30] + "..." if len(element.get('text', '')) > 30 else element.get('text', '')
                    
                    if element_type == 'text_frame_paragraph':
                        element['paragraph'].text = translated_text
                        updated_count += 1
                        print_info(f"Updated text_frame_paragraph: {original_text} -> {translated_text[:30]}...")
                        
                    elif element_type == 'table_cell':
                        element['cell'].text = translated_text
                        updated_count += 1
                        print_info(f"Updated table_cell: {original_text} -> {translated_text[:30]}...")
                        
                    elif element_type in ['chart_title', 'chart_axis_title']:
                        if element_type == 'chart_title':
                            element['chart'].chart_title.text_frame.text = translated_text
                        updated_count += 1
                        print_info(f"Updated {element_type}: {original_text} -> {translated_text[:30]}...")
                        
                    elif element_type == 'placeholder':
                        element['shape'].text = translated_text
                        updated_count += 1
                        print_info(f"Updated placeholder: {original_text} -> {translated_text[:30]}...")
                        
                    elif element_type == 'text_shape':
                        # Handle text shapes that aren't placeholders
                        if 'shape' in element:
                            element['shape'].text = translated_text
                            updated_count += 1
                            print_info(f"Updated text_shape: {original_text} -> {translated_text[:30]}...")
                        
                    elif element_type == 'layout_placeholder':
                        # Handle layout placeholders
                        if 'placeholder' in element:
                            element['placeholder'].text_frame.text = translated_text
                            updated_count += 1
                            print_info(f"Updated layout_placeholder: {original_text} -> {translated_text[:30]}...")
                        
                    elif element_type == 'notes_paragraph':
                        element['paragraph'].text = translated_text
                        updated_count += 1
                        print_info(f"Updated notes_paragraph: {original_text} -> {translated_text[:30]}...")
                        
                    elif element.get('requires_com', False):
                        # Store COM updates for later processing
                        com_updates_needed.append(element)
                        
                    elif element.get('requires_xml', False):
                        # XML updates would require more complex handling
                        print_warning(f"XML element update not yet implemented: {element_type}")
                        
                    else:
                        print_warning(f"Unknown element type for translation application: {element_type}")
                        
                except Exception as e:
                    print_warning(f"Failed to apply translation to {element.get('type', 'unknown')}: {e}")
                    import traceback
                    print_warning(f"Error details: {traceback.format_exc()}")
            
            # Save presentation with basic updates
            prs.save(output_path)
            print_success(f"Applied {updated_count} translations using python-pptx")
            
            # Apply COM updates if needed and available
            if com_updates_needed and HAS_COM and os.name == 'nt':
                print_info(f"Applying {len(com_updates_needed)} COM-based translations")
                self._apply_com_updates(com_updates_needed, output_path)
            elif com_updates_needed:
                print_warning(f"Skipped {len(com_updates_needed)} advanced elements (COM not available)")
            
            print_success(f"Advanced translation completed: {output_path}")
            return output_path
            
        except Exception as e:
            print_error(f"Failed to apply translations: {e}")
            return None

    def _apply_com_updates(self, com_elements: List[Dict], file_path: str):
        """Apply COM-based updates to advanced elements"""
        try:
            import win32com.client
            
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            ppt_app.Visible = False
            
            presentation = ppt_app.Presentations.Open(os.path.abspath(file_path))
            
            for element in com_elements:
                try:
                    slide_idx, shape_idx = element['location'][:2]
                    slide = presentation.Slides[slide_idx + 1]  # COM is 1-indexed
                    shape = slide.Shapes[shape_idx + 1]
                    
                    if element['type'] == 'smartart_node':
                        node_idx = element['node_index']
                        node = shape.SmartArt.AllNodes[node_idx + 1]
                        node.TextFrame2.TextRange.Text = element['translated_text']
                        
                    elif element['type'] == 'wordart':
                        shape.TextEffect.Text = element['translated_text']
                        
                    # OLE objects would need more specific handling per type
                        
                    print_success(f"Updated {element['type']} via COM")
                    
                except Exception as e:
                    print_warning(f"COM update failed for {element['type']}: {e}")
            
            presentation.Save()
            presentation.Close()
            ppt_app.Quit()
            
        except Exception as e:
            print_warning(f"COM updates failed: {e}")

# PDF processing removed to streamline the application

    def _analyze_pdf_type(self, file_path: str) -> str:
        """Analyze PDF to determine if it's text-based, scanned, or mixed content"""
        try:
            doc = fitz.open(file_path)
            total_text_chars = 0
            total_images = 0
            total_pages = len(doc)
            
            # Sample first 5 pages for analysis
            sample_pages = min(5, total_pages)
            
            for page_num in range(sample_pages):
                page = doc[page_num]
                
                # Count text characters
                text = page.get_text()
                total_text_chars += len(text.strip())
                
                # Count images
                image_list = page.get_images()
                total_images += len(image_list)
            
            doc.close()
            
            # Decision logic
            avg_text_per_page = total_text_chars / sample_pages
            avg_images_per_page = total_images / sample_pages
            
            if avg_text_per_page < 50 and avg_images_per_page > 0:
                return "scanned"  # Likely scanned document
            elif avg_text_per_page > 200 and avg_images_per_page > 0.5:
                return "mixed"    # Mixed content (text + images)
            else:
                return "text"     # Primarily text-based
                
        except Exception as e:
            print_warning(f"PDF analysis failed, defaulting to hybrid approach: {e}")
            return "mixed"

    def _process_with_hybrid_approach(self, file_path: str, target_lang: str, api_client: Any, progress: Optional[Any]) -> Optional[str]:
        """Comprehensive approach combining text extraction and OCR"""
        
        filename = os.path.basename(file_path)
        print_info(f"🔄 Using hybrid approach for comprehensive PDF processing")
        
        # Step 1: Extract text-based content
        text_elements = self._extract_text_based_content(file_path, progress)
        
        # Step 2: Process images/scanned content with OCR
        if HAS_IMAGE_LIBS and api_client:
            try:
                print_info("🔍 Processing images and scanned content with Gemini OCR...")
                ocr_elements = self._extract_with_gemini_ocr(file_path, api_client, progress)
                text_elements.extend(ocr_elements)
                print_info(f"📝 Gemini OCR found {len(ocr_elements)} additional text elements")
            except Exception as e:
                print_warning(f"OCR processing failed: {e}")
                
        # Step 3: Process translations
        if text_elements:
            print_info(f"🔄 Processing {len(text_elements)} total extracted elements")
            translated_elements = self._translate_pdf_elements(text_elements, target_lang, api_client, progress)
            return self._create_translated_pdf(file_path, translated_elements, progress)
            
        print_warning("No elements extracted by any engine")
        return None

    def _process_text_based_pdf(self, file_path: str, target_lang: str, api_client: Any, progress: Optional[Any]) -> Optional[str]:
        """Process text-based PDF using direct text extraction"""
        print_info("📄 Processing text-based PDF")
        
        elements = self._extract_text_based_content(file_path, progress)
        if elements:
            translated_elements = self._translate_pdf_elements(elements, target_lang, api_client, progress)
            return self._create_translated_pdf(file_path, translated_elements, progress)
        return None

    def _process_scanned_pdf_with_ocr(self, file_path: str, target_lang: str, api_client: Any, progress: Optional[Any]) -> Optional[str]:
        """Process scanned PDF using Gemini Vision API OCR"""
        if not HAS_IMAGE_LIBS or not api_client:
            raise Exception("OCR processing requires image libraries and API client")
            
        print_info("🔍 Processing scanned PDF with Gemini Vision OCR")
        
        elements = self._extract_with_gemini_ocr(file_path, api_client, progress)
        if elements:
            translated_elements = self._translate_pdf_elements(elements, target_lang, api_client, progress)
            return self._create_translated_pdf(file_path, translated_elements, progress)
        return None

    def _process_mixed_content_pdf(self, file_path: str, target_lang: str, api_client: Any, progress: Optional[Any]) -> Optional[str]:
        """Process mixed content PDF with intelligent content detection"""
        print_info("🔀 Processing mixed content PDF")
        return self._process_with_hybrid_approach(file_path, target_lang, api_client, progress)

    def _extract_text_based_content(self, file_path: str, progress: Optional[Any]) -> List[Dict]:
        """Extract text content from PDF using PyMuPDF and pdfplumber"""
        elements = []
        
        try:
            # Method 1: PyMuPDF for structured extraction
            print_info("🔤 Extracting text with PyMuPDF...")
            doc = fitz.open(file_path)
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                if progress:
                    progress.update(progress.task_ids[0] if progress.task_ids else None,
                                   description=f"[green]Extracting text from page {page_num+1}/{len(doc)}")
                
                # Extract text blocks with position information
                text_dict = page.get_text("dict")
                
                for block in text_dict["blocks"]:
                    if "lines" in block:  # Text block
                        for line in block["lines"]:
                            for span in line["spans"]:
                                text = span["text"].strip()
                                if text and should_translate(text):
                                    elements.append({
                                        'type': 'text_block',
                                        'page': page_num,
                                        'text': clean_text(text),
                                        'bbox': span["bbox"],  # Bounding box
                                        'font': span.get("font", ""),
                                        'size': span.get("size", 12),
                                        'extraction_method': 'pymupdf'
                                    })
            
            doc.close()
            
            # Method 2: pdfplumber for table detection
            print_info("📊 Extracting tables with pdfplumber...")
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    # Extract tables
                    tables = page.extract_tables()
                    for table_idx, table in enumerate(tables):
                        for row_idx, row in enumerate(table):
                            for col_idx, cell in enumerate(row):
                                if cell and should_translate(str(cell)):
                                    elements.append({
                                        'type': 'table_cell',
                                        'page': page_num,
                                        'text': clean_text(str(cell)),
                                        'table_idx': table_idx,
                                        'row': row_idx,
                                        'col': col_idx,
                                        'extraction_method': 'pdfplumber'
                                    })
                                    
        except Exception as e:
            print_warning(f"Text extraction failed: {e}")
            
        print_info(f"📄 Extracted {len(elements)} text elements")
        return elements

    def _extract_with_gemini_ocr(self, file_path: str, api_client: Any, progress: Optional[Any]) -> List[Dict]:
        """Extract text from PDF images using Gemini Vision API"""
        elements = []
        
        if not HAS_IMAGE_LIBS:
            print_warning("Image processing libraries not available for OCR")
            return elements
            
        try:
            print_info("🖼️ Converting PDF pages to images for OCR...")
            
            # Convert PDF to images
            images = convert_from_path(file_path, dpi=200, fmt='PNG')
            
            for page_num, image in enumerate(images):
                if progress:
                    progress.update(progress.task_ids[0] if progress.task_ids else None,
                                   description=f"[yellow]OCR processing page {page_num+1}/{len(images)}")
                
                # Convert PIL image to base64 for Gemini
                buffered = io.BytesIO()
                image.save(buffered, format="PNG")
                image_base64 = base64.b64encode(buffered.getvalue()).decode()
                
                # Use Gemini Vision API for OCR
                ocr_text = self._gemini_vision_ocr(image_base64, api_client)
                
                if ocr_text and should_translate(ocr_text):
                    # Split OCR text into logical chunks
                    text_chunks = self._split_ocr_text(ocr_text)
                    
                    for chunk_idx, chunk in enumerate(text_chunks):
                        if chunk.strip() and should_translate(chunk):
                            elements.append({
                                'type': 'ocr_text',
                                'page': page_num,
                                'text': clean_text(chunk),
                                'chunk_idx': chunk_idx,
                                'extraction_method': 'gemini_ocr'
                            })
                            
        except Exception as e:
            print_warning(f"Gemini OCR processing failed: {e}")
            
        print_info(f"🔍 OCR extracted {len(elements)} text elements")
        return elements

    def _gemini_vision_ocr(self, image_base64: str, api_client: Any) -> str:
        """Use Gemini Vision API to perform OCR on image"""
        try:
            # Prepare the vision prompt for OCR
            ocr_prompt = """Extract all text from this image. Please:
1. Maintain the original text layout and structure
2. Preserve line breaks and paragraph separation
3. Include all visible text including headers, body text, captions, etc.
4. Do not add any explanations or comments
5. Return only the extracted text content"""

            # Call Gemini Vision API
            response = api_client.chat.completions.create(
                model="gemini-2.0-flash",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": ocr_prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{image_base64}"
                                }
                            }
                        ]
                    }
                ]
            )
            
            # Extract text from response
            if hasattr(response, 'choices') and response.choices:
                return response.choices[0].message.content
            else:
                # Try alternative response formats
                response_text = str(response)
                return response_text if len(response_text) > 20 else ""
                
        except Exception as e:
            print_warning(f"Gemini Vision OCR failed: {e}")
            return ""

    def _split_ocr_text(self, text: str) -> List[str]:
        """Split OCR text into logical chunks for better translation"""
        # Split by double line breaks (paragraphs)
        paragraphs = text.split('\n\n')
        
        chunks = []
        for para in paragraphs:
            # Further split long paragraphs by sentences
            if len(para) > 300:
                sentences = para.split('. ')
                current_chunk = ""
                
                for sentence in sentences:
                    if len(current_chunk + sentence) < 200:
                        current_chunk += sentence + ". "
                    else:
                        if current_chunk.strip():
                            chunks.append(current_chunk.strip())
                        current_chunk = sentence + ". "
                
                if current_chunk.strip():
                    chunks.append(current_chunk.strip())
            else:
                if para.strip():
                    chunks.append(para.strip())
        
        return chunks

    def _translate_pdf_elements(self, elements: List[Dict], target_lang: str, api_client: Any, progress: Optional[Any]) -> List[Dict]:
        """Translate extracted PDF elements"""
        if not elements:
            return []
            
        print_info(f"🔄 Translating {len(elements)} PDF elements")
        
        # Group texts for batch translation
        texts_to_translate = []
        element_map = {}
        
        for idx, element in enumerate(elements):
            if element['text'] and should_translate(element['text']):
                clean_text_value = clean_text(element['text'])
                texts_to_translate.append(clean_text_value)
                element_map[len(texts_to_translate) - 1] = idx
        
        if not texts_to_translate:
            return elements
            
        print_info(f"📝 Processing {len(texts_to_translate)} texts for translation")
        
        # Use existing translation function
        translated_texts = translate_batch(texts_to_translate, target_lang, api_client)
        
        # Map translations back to elements
        for trans_idx, translated_text in enumerate(translated_texts):
            if trans_idx in element_map:
                original_idx = element_map[trans_idx]
                elements[original_idx]['translated_text'] = translated_text
                
        return elements

    def _create_translated_pdf(self, file_path: str, translated_elements: List[Dict], progress: Optional[Any]) -> Optional[str]:
        """Create translated PDF with preserved layout"""
        
        # Create output path
        filename = os.path.basename(file_path)
        base_name, ext = os.path.splitext(filename)
        safe_base_name = re.sub(r'[^\w\-]', '_', base_name)
        safe_base_name = re.sub(r'_+', '_', safe_base_name).strip('_')
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        project_dir = os.path.dirname(os.path.abspath(__file__))
        output_dir = os.path.join(project_dir, "output")
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f"TRANSLATED_{safe_base_name}_{timestamp}{ext}")
        
        try:
            print_info(f"📄 Creating translated PDF: {output_path}")
            
            # Create new PDF with translations
            doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Group elements by page
            pages_content = {}
            for element in translated_elements:
                if 'translated_text' in element:
                    page_num = element.get('page', 0)
                    if page_num not in pages_content:
                        pages_content[page_num] = []
                    pages_content[page_num].append(element)
            
            # Build PDF content
            for page_num in sorted(pages_content.keys()):
                # Add page header
                story.append(Paragraph(f"<b>Page {page_num + 1}</b>", styles['Heading2']))
                story.append(Spacer(1, 12))
                
                # Add translated content
                for element in pages_content[page_num]:
                    translated_text = element['translated_text']
                    
                    # Format based on element type
                    if element['type'] == 'table_cell':
                        text_content = f"<i>Table [{element.get('row', 0)}, {element.get('col', 0)}]:</i> {translated_text}"
                    else:
                        text_content = translated_text
                    
                    story.append(Paragraph(text_content, styles['Normal']))
                    story.append(Spacer(1, 6))
                
                # Add page break except for last page
                if page_num < max(pages_content.keys()):
                    story.append(Spacer(1, 20))
            
            # Build the PDF
            doc.build(story)
            print_success(f"✅ Translated PDF created successfully: {output_path}")
            
            return output_path
            
        except Exception as e:
            print_error(f"❌ Failed to create translated PDF: {e}")
            return None

def process_powerpoint_file_basic(
    input_path: str, 
    target_lang: str = "ja", 
    api_client: Any = None,
    progress: Optional[Any] = None
) -> Optional[str]:
    """
    Process PowerPoint file using basic python-pptx approach: read, translate and save with original format.
    
    Args:
        input_path: Path to the PowerPoint file
        target_lang: Target language code
        api_client: OpenAI client instance
        progress: Optional rich progress instance
        
    Returns:
        Optional[str]: Path to the translated file, or None if processing failed
    """
    try:
        # Import python-pptx here to ensure it's loaded
        from pptx import Presentation
        
        # Create output file path with completely safe filename
        filename = os.path.basename(input_path)
        base_name, ext = os.path.splitext(filename)
        # Create completely safe filename with timestamp
        import re
        # Remove all special characters and replace with underscores
        safe_base_name = re.sub(r'[^\w\-]', '_', base_name)
        # Remove multiple consecutive underscores
        safe_base_name = re.sub(r'_+', '_', safe_base_name).strip('_')
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")

        # Create output directory at the same level as the script
        project_dir = os.path.dirname(os.path.abspath(__file__))
        output_dir = os.path.join(project_dir, "output")
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f"TRANSLATED_{safe_base_name}_{timestamp}{ext}")

        print_info(f"Processing PowerPoint file: {filename}")

        # Task IDs for progress tracking
        scan_task = None
        file_task = None
        
        if progress:
            file_task = progress.add_task(f"[cyan]Processing {filename}", total=1.0)
            scan_task = progress.add_task("[green]Scanning slides", total=1.0)

        # Load the presentation
        prs = Presentation(input_path)
        
        if progress:
            # Update progress: file opened
            progress.update(file_task, completed=0.2, description=f"[cyan]Processing {filename} - File opened")
            # Update scanning task: starting scan
            progress.update(scan_task, completed=0.1, description="[green]Scanning slides")
        
        # Collect all text elements that need translation
        texts_to_translate = []
        text_references = []
        
        # Process each slide
        slide_count = len(prs.slides)
        
        for slide_idx, slide in enumerate(prs.slides):
            if progress:
                progress.update(scan_task, completed=0.1 + (0.8 * ((slide_idx + 1) / slide_count)), 
                               description=f"[green]Scanning slide {slide_idx+1}/{slide_count}")
            
            print_info(f"Scanning slide {slide_idx + 1}/{slide_count}")
            
            # Process shapes in the slide (text boxes, titles, etc.)
            for shape_idx, shape in enumerate(slide.shapes):
                # Process text frames
                if hasattr(shape, "text") and shape.text:
                    text = shape.text
                    if should_translate(text):
                        clean_text_value = clean_text(text)
                        texts_to_translate.append(clean_text_value)
                        # Store reference as (slide index, shape index)
                        text_references.append((slide_idx, shape_idx))
                        preview = clean_text_value[:30] + "..." if len(clean_text_value) > 30 else clean_text_value
                        print_info(f"Found text in slide {slide_idx + 1}, shape {shape_idx + 1}: {preview}")
                
                # Process text in tables
                if hasattr(shape, "has_table") and shape.has_table:
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            if cell.text and should_translate(cell.text):
                                clean_text_value = clean_text(cell.text)
                                texts_to_translate.append(clean_text_value)
                                # Store reference as (slide index, shape index, "table", row index, col index)
                                text_references.append((slide_idx, shape_idx, "table", row_idx, col_idx))
                                preview = clean_text_value[:30] + "..." if len(clean_text_value) > 30 else clean_text_value
                                print_info(f"Found text in slide {slide_idx + 1}, table cell [{row_idx + 1}, {col_idx + 1}]: {preview}")
                
                # Process text in groups (handle nested shapes)
                if hasattr(shape, "group_items") and shape.group_items:
                    for group_idx, group_shape in enumerate(shape.group_items):
                        if hasattr(group_shape, "text") and group_shape.text:
                            text = group_shape.text
                            if should_translate(text):
                                clean_text_value = clean_text(text)
                                texts_to_translate.append(clean_text_value)
                                # Store reference as (slide index, shape index, "group", group shape index)
                                text_references.append((slide_idx, shape_idx, "group", group_idx))
                                preview = clean_text_value[:30] + "..." if len(clean_text_value) > 30 else clean_text_value
                                print_info(f"Found text in slide {slide_idx + 1}, grouped shape {group_idx + 1}: {preview}")

        # Process notes
        for slide_idx, slide in enumerate(prs.slides):
            if hasattr(slide, "notes_slide") and slide.notes_slide and hasattr(slide.notes_slide, "notes_text_frame"):
                text = slide.notes_slide.notes_text_frame.text
                if text and should_translate(text):
                    clean_text_value = clean_text(text)
                    texts_to_translate.append(clean_text_value)
                    # Store reference as (slide index, "notes")
                    text_references.append((slide_idx, "notes"))
                    preview = clean_text_value[:30] + "..." if len(clean_text_value) > 30 else clean_text_value
                    print_info(f"Found notes in slide {slide_idx + 1}: {preview}")
        
        # Finish scanning phase
        if progress:
            progress.update(scan_task, completed=1.0, description="[green]Scanning complete")
        
        # Skip if no text to translate
        if not texts_to_translate:
            print_success(f"No text to translate in presentation '{filename}'.")
            if progress:
                progress.update(file_task, completed=1.0, description=f"[cyan]Processing {filename} - No translation needed")
            # Still save a copy in the output directory
            prs.save(output_path)
            print_success(f"File saved at: {output_path}")
            return output_path
        
        # Add translation task
        total_elements = len(texts_to_translate)
        if progress:
            trans_task = progress.add_task("[yellow]Translating content", total=total_elements)
            progress.update(file_task, completed=0.4, description=f"[cyan]Processing {filename} - Translation in progress")
            
        # Split into batches for processing
        batch_size = 50  # Smaller batch size for better progress tracking
        total_batches = (len(texts_to_translate) - 1) // batch_size + 1
        processed_count = 0
        
        print_info(f"Preparing to translate {len(texts_to_translate)} text segments in {total_batches} batches.")
        
        for i in range(0, len(texts_to_translate), batch_size):
            batch_texts = texts_to_translate[i:i+batch_size]
            batch_refs = text_references[i:i+batch_size]
            current_batch_num = i // batch_size + 1
            
            print_info(f"Translating batch {current_batch_num}/{total_batches} ({len(batch_texts)} texts)")
            
            # Define a progress callback for translation
            def update_translation_progress(current: int, total: int) -> None:
                if progress:
                    progress.update(trans_task, 
                                   description=f"[yellow]Translating batch {current_batch_num}/{total_batches} (Attempt {current+1}/{total+1})")
            
            # Translate batch
            translated_batch = translate_batch(
                batch_texts, 
                target_lang,
                api_client, 
                progress_callback=update_translation_progress
            )
            
            # Update translated content
            print_info(f"Updating content for batch {current_batch_num}...")
            
            for j, ref in enumerate(batch_refs):
                # Update progress
                if progress:
                    processed_count += 1
                    progress.update(trans_task, completed=processed_count)
                    
                if j < len(translated_batch) and translated_batch[j]:
                    try:
                        # Regular shape
                        if len(ref) == 2 and isinstance(ref[0], int) and isinstance(ref[1], int):
                            slide_idx, shape_idx = ref
                            prs.slides[slide_idx].shapes[shape_idx].text = translated_batch[j]
                            print_success(f"Updated text in slide {slide_idx + 1}, shape {shape_idx + 1}")
                        
                        # Table cell
                        elif len(ref) == 5 and ref[2] == "table":
                            slide_idx, shape_idx, _, row_idx, col_idx = ref
                            prs.slides[slide_idx].shapes[shape_idx].table.rows[row_idx].cells[col_idx].text = translated_batch[j]
                            print_success(f"Updated text in slide {slide_idx + 1}, table cell [{row_idx + 1}, {col_idx + 1}]")
                        
                        # Group shape
                        elif len(ref) == 4 and ref[2] == "group":
                            slide_idx, shape_idx, _, group_idx = ref
                            prs.slides[slide_idx].shapes[shape_idx].group_items[group_idx].text = translated_batch[j]
                            print_success(f"Updated text in slide {slide_idx + 1}, grouped shape {group_idx + 1}")
                        
                        # Notes
                        elif len(ref) == 2 and ref[1] == "notes":
                            slide_idx = ref[0]
                            if prs.slides[slide_idx].notes_slide and hasattr(prs.slides[slide_idx].notes_slide, "notes_text_frame"):
                                prs.slides[slide_idx].notes_slide.notes_text_frame.text = translated_batch[j]
                                print_success(f"Updated notes in slide {slide_idx + 1}")
                        
                        else:
                            print_warning(f"Unknown reference type: {ref}")
                            
                    except Exception as update_err:
                        print_warning(f"Error updating content: {str(update_err)}")
        
        # Update file task
        if progress:
            progress.update(file_task, completed=0.8, description=f"[cyan]Processing {filename} - Saving file")
            progress.update(trans_task, completed=total_elements, description="[yellow]Translation complete")
            
        # Save the translated presentation
        print_info(f"Saving translated PowerPoint to: {output_path}")
        prs.save(output_path)
        print_success(f"File saved successfully: {output_path}")
        
        if progress:
            progress.update(file_task, completed=1.0, description=f"[cyan]Processing {filename} - Complete")
            
        return output_path
        
    except Exception as e:
        print_error(f"Error processing PowerPoint file '{input_path}': {str(e)}")
        if progress and file_task:
            progress.update(file_task, completed=1.0, description=f"[red]Processing {filename} - Failed")
        return None

def process_word_file(
    input_path: str, 
    target_lang: str = "ja", 
    api_client: Any = None,
    progress: Optional[Any] = None
) -> Optional[str]:
    """
    Process Word file: read, translate and save with original format.
    
    Args:
        input_path: Path to the Word file
        target_lang: Target language code
        api_client: OpenAI client instance
        progress: Optional rich progress instance
        
    Returns:
        Optional[str]: Path to the translated file, or None if processing failed
    """
    try:
        # Import python-docx here to ensure it's loaded
        from docx import Document
        
        # Create output file path
        filename = os.path.basename(input_path)
        base_name, ext = os.path.splitext(filename)

        # Create output directory at the same level as the script
        project_dir = os.path.dirname(os.path.abspath(__file__))
        output_dir = os.path.join(project_dir, "output")
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f"{base_name}-translated{ext}")

        print_info(f"Processing Word file: {filename}")

        # Task IDs for progress tracking
        scan_task = None
        file_task = None
        
        if progress:
            file_task = progress.add_task(f"[cyan]Processing {filename}", total=1.0)
            scan_task = progress.add_task("[green]Scanning document", total=1.0)

        # Load the document
        doc = Document(input_path)
        
        if progress:
            # Update progress: file opened
            progress.update(file_task, completed=0.2, description=f"[cyan]Processing {filename} - File opened")
            # Update scanning task: starting scan
            progress.update(scan_task, completed=0.1, description="[green]Scanning document content")
        
        # Collect all text elements that need translation
        texts_to_translate = []
        text_references = []
        
        # Process paragraphs
        para_count = len(doc.paragraphs)
        for para_idx, para in enumerate(doc.paragraphs):
            if progress:
                scan_progress = 0.1 + (0.3 * ((para_idx + 1) / para_count))
                progress.update(scan_task, completed=scan_progress, 
                               description=f"[green]Scanning paragraphs {para_idx+1}/{para_count}")
                
            if para.text and should_translate(para.text):
                clean_text_value = clean_text(para.text)
                texts_to_translate.append(clean_text_value)
                # Store reference as ("paragraph", paragraph index)
                text_references.append(("paragraph", para_idx))
                preview = clean_text_value[:30] + "..." if len(clean_text_value) > 30 else clean_text_value
                print_info(f"Found text in paragraph {para_idx + 1}: {preview}")
        
        # Process tables
        table_count = len(doc.tables)
        for table_idx, table in enumerate(doc.tables):
            if progress:
                scan_progress = 0.4 + (0.3 * ((table_idx + 1) / (table_count or 1)))
                progress.update(scan_task, completed=scan_progress, 
                               description=f"[green]Scanning tables {table_idx+1}/{table_count}")
                
            for row_idx, row in enumerate(table.rows):
                for cell_idx, cell in enumerate(row.cells):
                    # Check cell paragraphs
                    for para_idx, para in enumerate(cell.paragraphs):
                        if para.text and should_translate(para.text):
                            clean_text_value = clean_text(para.text)
                            texts_to_translate.append(clean_text_value)
                            # Store reference as ("table", table index, row index, cell index, paragraph index)
                            text_references.append(("table", table_idx, row_idx, cell_idx, para_idx))
                            preview = clean_text_value[:30] + "..." if len(clean_text_value) > 30 else clean_text_value
                            print_info(f"Found text in table {table_idx + 1}, cell [{row_idx + 1}, {cell_idx + 1}]: {preview}")
        
        # Process headers and footers
        section_count = len(doc.sections)
        for section_idx, section in enumerate(doc.sections):
            if progress:
                scan_progress = 0.7 + (0.25 * ((section_idx + 1) / (section_count or 1)))
                progress.update(scan_task, completed=scan_progress, 
                               description=f"[green]Scanning sections {section_idx+1}/{section_count}")
                
            # Process headers
            for header_type in ['header', 'first_page_header', 'even_page_header']:
                header = getattr(section, header_type)
                if header:
                    for para_idx, para in enumerate(header.paragraphs):
                        if para.text and should_translate(para.text):
                            clean_text_value = clean_text(para.text)
                            texts_to_translate.append(clean_text_value)
                            # Store reference as ("header", section index, header type, paragraph index)
                            text_references.append(("header", section_idx, header_type, para_idx))
                            preview = clean_text_value[:30] + "..." if len(clean_text_value) > 30 else clean_text_value
                            print_info(f"Found text in section {section_idx + 1}, {header_type}, paragraph {para_idx + 1}: {preview}")
            
            # Process footers
            for footer_type in ['footer', 'first_page_footer', 'even_page_footer']:
                footer = getattr(section, footer_type)
                if footer:
                    for para_idx, para in enumerate(footer.paragraphs):
                        if para.text and should_translate(para.text):
                            clean_text_value = clean_text(para.text)
                            texts_to_translate.append(clean_text_value)
                            # Store reference as ("footer", section index, footer type, paragraph index)
                            text_references.append(("footer", section_idx, footer_type, para_idx))
                            preview = clean_text_value[:30] + "..." if len(clean_text_value) > 30 else clean_text_value
                            print_info(f"Found text in section {section_idx + 1}, {footer_type}, paragraph {para_idx + 1}: {preview}")
        
        # Finish scanning phase
        if progress:
            progress.update(scan_task, completed=1.0, description="[green]Scanning complete")
        
        # Skip if no text to translate
        if not texts_to_translate:
            print_success(f"No text to translate in document '{filename}'.")
            if progress:
                progress.update(file_task, completed=1.0, description=f"[cyan]Processing {filename} - No translation needed")
            # Still save a copy in the output directory
            doc.save(output_path)
            print_success(f"File saved at: {output_path}")
            return output_path
        
        # Add translation task
        total_elements = len(texts_to_translate)
        if progress:
            trans_task = progress.add_task("[yellow]Translating content", total=total_elements)
            progress.update(file_task, completed=0.4, description=f"[cyan]Processing {filename} - Translation in progress")
            
        # Split into batches for processing
        batch_size = 50  # Smaller batch size for better progress tracking
        total_batches = (len(texts_to_translate) - 1) // batch_size + 1
        processed_count = 0
        
        print_info(f"Preparing to translate {len(texts_to_translate)} text segments in {total_batches} batches.")
        
        for i in range(0, len(texts_to_translate), batch_size):
            batch_texts = texts_to_translate[i:i+batch_size]
            batch_refs = text_references[i:i+batch_size]
            current_batch_num = i // batch_size + 1
            
            print_info(f"Translating batch {current_batch_num}/{total_batches} ({len(batch_texts)} texts)")
            
            # Define a progress callback for translation
            def update_translation_progress(current: int, total: int) -> None:
                if progress:
                    progress.update(trans_task, 
                                   description=f"[yellow]Translating batch {current_batch_num}/{total_batches} (Attempt {current+1}/{total+1})")
            
            # Translate batch
            translated_batch = translate_batch(
                batch_texts, 
                target_lang,
                api_client, 
                progress_callback=update_translation_progress
            )
            
            # Update translated content
            print_info(f"Updating content for batch {current_batch_num}...")
            
            for j, ref in enumerate(batch_refs):
                # Update progress
                if progress:
                    processed_count += 1
                    progress.update(trans_task, completed=processed_count)
                    
                if j < len(translated_batch) and translated_batch[j]:
                    try:
                        # Paragraph
                        if ref[0] == "paragraph":
                            _, para_idx = ref
                            doc.paragraphs[para_idx].text = translated_batch[j]
                            print_success(f"Updated text in paragraph {para_idx + 1}")
                        
                        # Table cell
                        elif ref[0] == "table":
                            _, table_idx, row_idx, cell_idx, para_idx = ref
                            doc.tables[table_idx].rows[row_idx].cells[cell_idx].paragraphs[para_idx].text = translated_batch[j]
                            print_success(f"Updated text in table {table_idx + 1}, cell [{row_idx + 1}, {cell_idx + 1}]")
                        
                        # Header
                        elif ref[0] == "header":
                            _, section_idx, header_type, para_idx = ref
                            header = getattr(doc.sections[section_idx], header_type)
                            header.paragraphs[para_idx].text = translated_batch[j]
                            print_success(f"Updated text in section {section_idx + 1}, {header_type}")
                        
                        # Footer
                        elif ref[0] == "footer":
                            _, section_idx, footer_type, para_idx = ref
                            footer = getattr(doc.sections[section_idx], footer_type)
                            footer.paragraphs[para_idx].text = translated_batch[j]
                            print_success(f"Updated text in section {section_idx + 1}, {footer_type}")
                        
                        else:
                            print_warning(f"Unknown reference type: {ref[0]}")
                            
                    except Exception as update_err:
                        print_warning(f"Error updating content: {str(update_err)}")
        
        # Update file task
        if progress:
            progress.update(file_task, completed=0.8, description=f"[cyan]Processing {filename} - Saving file")
            progress.update(trans_task, completed=total_elements, description="[yellow]Translation complete")
            
        # Save the translated document
        print_info(f"Saving translated Word document to: {output_path}")
        doc.save(output_path)
        print_success(f"File saved successfully: {output_path}")
        
        if progress:
            progress.update(file_task, completed=1.0, description=f"[cyan]Processing {filename} - Complete")
            
        return output_path
        
    except Exception as e:
        print_error(f"Error processing Word file '{input_path}': {str(e)}")
        if progress and file_task:
            progress.update(file_task, completed=1.0, description=f"[red]Processing {filename} - Failed")
        return None 

def process_file(
    input_path: str, 
    target_lang: str = "ja",
    progress: Optional[Any] = None
) -> Optional[str]:
    """
    Process a file based on its type.
    
    Args:
        input_path: Path to the file
        target_lang: Target language code
        progress: Optional rich progress instance
        
    Returns:
        Optional[str]: Path to the translated file, or None if processing failed
    """
    # Check if file exists
    if not os.path.exists(input_path):
        print_error(f"File does not exist: {input_path}")
        return None
    
    # Import required libraries
    try:
        from dotenv import load_dotenv
        from openai import OpenAI
    except ImportError as e:
        print_error(f"Failed to import required libraries: {str(e)}")
        return None
    
    # Load environment variables
    load_dotenv()
    
    # Check if API key is available
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        print_error("GEMINI_API_KEY environment variable not found.")
        print_info("Please add your API key to the .env file and try again.")
        return None
    
    # Initialize API client
    try:
        client = OpenAI(
            base_url="https://generativelanguage.googleapis.com/v1beta/",
            api_key=api_key,
        )
    except Exception as e:
        print_error(f"Failed to initialize API client: {str(e)}")
        return None
    
    # Determine file type
    file_type = get_file_type(input_path)
    
    # Process file based on type
    if file_type == DocType.EXCEL:
        print_info(f"Processing Excel file: {os.path.basename(input_path)}")
        return process_excel_file(input_path, target_lang, client, progress)
    elif file_type == DocType.WORD:
        print_info(f"Processing Word file: {os.path.basename(input_path)}")
        return process_word_file(input_path, target_lang, client, progress)
    elif file_type == DocType.POWERPOINT:
        print_info(f"Processing PowerPoint file: {os.path.basename(input_path)}")
        return process_powerpoint_file(input_path, target_lang, client, progress)
    # PDF processing removed
    else:
        print_error(f"Unsupported file type: {os.path.basename(input_path)}")
        return None

def process_directory(
    input_dir: str, 
    target_lang: str = "ja"
) -> Tuple[List[str], List[str]]:
    """
    Process all supported files in the input directory.
    
    Args:
        input_dir: Path to the input directory
        target_lang: Target language code
        
    Returns:
        Tuple[List[str], List[str]]: Lists of successful and failed files
    """
    # Ensure directory path exists
    if not os.path.isdir(input_dir):
        print_error(f"Directory does not exist: {input_dir}")
        return [], []

    # Find all supported files in the directory
    excel_files = glob.glob(os.path.join(input_dir, "*.xlsx")) + glob.glob(os.path.join(input_dir, "*.xls"))
    word_files = glob.glob(os.path.join(input_dir, "*.docx")) + glob.glob(os.path.join(input_dir, "*.doc"))
    ppt_files = glob.glob(os.path.join(input_dir, "*.pptx")) + glob.glob(os.path.join(input_dir, "*.ppt"))
    
    all_files = excel_files + word_files + ppt_files

    if not all_files:
        print_warning(f"No supported files found in directory: {input_dir}")
        print_info("Supported formats: .xlsx, .xls, .docx, .doc, .pptx, .ppt")
        return [], []

    print_info(f"Found {len(excel_files)} Excel files, {len(word_files)} Word files, and {len(ppt_files)} PowerPoint files")

    # Initialize rich progress if available
    progress = None
    if HAS_RICH:
        progress = Progress(
            SpinnerColumn(),
            TextColumn("[progress.description]{task.description}"),
            BarColumn(),
            TextColumn("[progress.percentage]{task.percentage:>3.0f}%"),
            console=console
        )
    else:
        print_info("Install 'rich' package for enhanced progress display")
    
    # Process each file
    successful_files = []
    failed_files = []
    
    # Start progress tracking if available
    if progress:
        with progress:
            main_task = progress.add_task(f"[blue]Processing {len(all_files)} files", total=len(all_files))
            
            for file_idx, file_path in enumerate(all_files):
                # Skip temporary files (usually starting with ~$)
                if os.path.basename(file_path).startswith('~$'):
                    print_info(f"Skipping temporary file: {os.path.basename(file_path)}")
                    continue

                progress.update(main_task, description=f"[blue]Processing file {file_idx+1}/{len(all_files)}")
                
                output_file = process_file(file_path, target_lang, progress)
                
                if output_file:
                    successful_files.append(os.path.basename(file_path))
                else:
                    failed_files.append(os.path.basename(file_path))
                    
                progress.update(main_task, advance=1)
    else:
        # Process without rich progress display
        for file_idx, file_path in enumerate(all_files):
            # Skip temporary files (usually starting with ~$)
            if os.path.basename(file_path).startswith('~$'):
                print_info(f"Skipping temporary file: {os.path.basename(file_path)}")
                continue

            print_info(f"Processing file {file_idx+1}/{len(all_files)}: {os.path.basename(file_path)}")
            
            output_file = process_file(file_path, target_lang)
            
            if output_file:
                successful_files.append(os.path.basename(file_path))
            else:
                failed_files.append(os.path.basename(file_path))

    print_header("Directory Processing Summary")
    print_success(f"Successfully processed: {len(successful_files)} files")
    if failed_files:
        print_warning(f"Failed to process: {len(failed_files)} files")
        if len(failed_files) <= 5:
            for failed in failed_files:
                print_warning(f"  - {failed}")
        else:
            for failed in failed_files[:5]:
                print_warning(f"  - {failed}")
            print_warning(f"  ... and {len(failed_files) - 5} more")

    return successful_files, failed_files

def process_powerpoint_file(
    input_path: str, 
    target_lang: str = "ja", 
    api_client: Any = None,
    progress: Optional[Any] = None
) -> Optional[str]:
    """
    PowerPoint file processor - using basic approach for reliability.
    
    The advanced processor with multi-engine approach has been temporarily disabled
    due to file corruption issues. Using the reliable basic processor instead.
    
    Args:
        input_path: Path to the PowerPoint file
        target_lang: Target language code
        api_client: OpenAI client instance
        progress: Optional rich progress instance
        
    Returns:
        Optional[str]: Path to the translated file, or None if processing failed
    """
    filename = os.path.basename(input_path)
    
    print_info(f"🚀 Processing PowerPoint file with basic processor: {filename}")
    
    try:
        # Use basic processor directly for reliability
        result = process_powerpoint_file_basic(input_path, target_lang, api_client, progress)
        
        if result:
            print_success(f"✅ PowerPoint processing completed successfully: {filename}")
            return result
        else:
            print_error(f"❌ PowerPoint processing failed for: {filename}")
            return None
            
    except Exception as e:
        print_error(f"❌ PowerPoint processing failed for {filename}: {str(e)}")
        return None

# PDF file processing removed to streamline the application

def main():
    """Main entry point for the translator"""
    # Display header
    print_header("Office Document Translator")
    
    # Setup argument parser
    parser = argparse.ArgumentParser(
        description='Translate Microsoft Office documents (Excel, Word, PowerPoint, PDF) between Japanese, English, and Vietnamese',
        formatter_class=argparse.ArgumentDefaultsHelpFormatter
    )
    
    # Add arguments
    parser.add_argument('--to', choices=['ja', 'vi', 'en', 'th', 'zh', 'ko'], default='ja',
                        help='Target language (ja: Japanese, vi: Vietnamese, en: English, th: Thai, zh: Chinese, ko: Korean)')
    parser.add_argument('--file', type=str, help='Path to a specific file to translate')
    parser.add_argument('--dir', type=str, help='Path to a directory containing files to translate')
    parser.add_argument('--output-dir', type=str, help='Path to the output directory (defaults to "./output")')
    parser.add_argument('--version', action='version', version='Office Document Translator v1.0.0')
    
    # Parse arguments
    args = parser.parse_args()
    
    # Check if required libraries are installed
    if not check_and_install_dependencies():
        print_error("Failed to install or load required dependencies.")
        return 1
    
    # Set output directory if specified
    # This part is a bit redundant as the individual processors also create the output dir
    # but it's good for clarity and potential future use.
    output_dir_arg = args.output_dir
    if output_dir_arg:
        if not os.path.exists(output_dir_arg):
            os.makedirs(output_dir_arg)
            print_info(f"Output directory created: {output_dir_arg}")
        # Note: Individual processors will use their own output logic within script_dir/output for now.
        # This global output_dir_arg is not directly passed to them yet.
        print_info(f"Global output directory specified: {output_dir_arg} (Note: processors use local ./output)")

    # Create input directory if not specified and no file/dir given
    if not args.file and not args.dir:
        script_dir = os.path.dirname(os.path.abspath(__file__))
        input_dir = os.path.join(script_dir, "input")
        
        if not os.path.exists(input_dir):
            os.makedirs(input_dir)
            print_info(f"Created 'input' directory at: {input_dir}")
            print_info("Please place files to translate in this directory or use --file/--dir arguments.")
            return 0 # Exit gracefully if input dir was just created
        
        args.dir = input_dir # Default to processing the input directory
        print_info(f"No file or directory specified, defaulting to input directory: {args.dir}")

    # Process based on provided arguments
    target_lang = args.to
    language_names = {
        'ja': 'Japanese',
        'vi': 'Vietnamese',
        'en': 'English',
        'th': 'Thai',
        'zh': 'Chinese (Simplified)',
        'ko': 'Korean'
    }
    print_info(f"Target language: {language_names.get(target_lang, target_lang)}")
    
    start_time = time.time()
    
    if args.file:
        # Process single file
        if os.path.exists(args.file):
            output_file = process_file(args.file, target_lang)
            if output_file:
                print_success(f"File successfully translated: {output_file}")
            else:
                print_error(f"Failed to translate file: {args.file}")
                return 1
        else:
            print_error(f"File not found: {args.file}")
            return 1
    elif args.dir:
        # Process directory
        if os.path.isdir(args.dir):
            successful, failed = process_directory(args.dir, target_lang)
            if not successful and failed:
                print_warning("No files were successfully translated.")
                # return 1 # Don't exit if some files failed but others might have succeeded
            elif not successful and not failed:
                print_info("No files found to process in the directory.")
        else:
            print_error(f"Directory not found: {args.dir}")
            return 1
    else:
        # This case should ideally be handled by the default input directory logic
        print_error("No file or directory specified, and default input directory is empty or not found.")
        return 1
    
    end_time = time.time()
    elapsed_time = end_time - start_time
    
    # Format time nicely
    if elapsed_time < 60:
        time_str = f"{elapsed_time:.2f} seconds"
    elif elapsed_time < 3600:
        minutes = int(elapsed_time // 60)
        seconds = elapsed_time % 60
        time_str = f"{minutes} minute{'s' if minutes != 1 else ''} and {seconds:.2f} seconds"
    else:
        hours = int(elapsed_time // 3600)
        minutes = int((elapsed_time % 3600) // 60)
        seconds = elapsed_time % 60
        time_str = f"{hours} hour{'s' if hours != 1 else ''}, {minutes} minute{'s' if minutes != 1 else ''}, and {seconds:.2f} seconds"
    
    print_info(f"Total execution time: {time_str}")
    
    return 0

if __name__ == "__main__":
    sys.exit(main()) 